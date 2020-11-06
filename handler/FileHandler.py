from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfinterp import resolve1
import docx2txt
from pptx import Presentation

import io
import pickle

class FileHandler():  
    def __init__(self):
        None
        
    def save_data(self, file_name, data):
        with open(file_name, "wb") as f:
            pickle.dump(data, f, -1)
        
    def load_data(self, file_name):
        with open(file_name, "rb") as f:
            data = pickle.load(f)
        return data
    
    def save_txt(self, file_name, line_list, encoding="utf-8", separator="\t"):
        f = open(file_name, "w", encoding=encoding)
        for line in line_list:
            if type(line) == "str":
                f.write(line.replace("\n", " "))
            else:
                new_line = ""
                for col in line:
                    new_line += str(col).replace("\n", " ") + separator
                f.write(new_line.strip())
            f.write("\n")
        f.close()
    
    def load_txt(self, file_name, encoding="utf-8", separator="\t"):
        line_list = []
        f = open(file_name, encoding=encoding)
        for line in f:
            line = line.replace("\n", "").split(separator)
            line_list.append(line)
        return line_list

    def pdf_converter(self, input_filename, output_filename):
        input_file = open(input_filename, "rb")
        pdf_parser = PDFParser(input_file)
        pdf_document = PDFDocument(pdf_parser)
        page_count = range(resolve1(pdf_document.catalog["Pages"])["Count"])
        string_io = io.StringIO()
        if not page_count:
            page_number_set = set()
        else:
            page_number_set = set(page_count)      
        resource_manager = PDFResourceManager() 
        converter = TextConverter(resource_manager, string_io, laparams=LAParams())     
        page_interpreter = PDFPageInterpreter(resource_manager, converter)
        for page in PDFPage.get_pages(input_file, page_number_set, caching=True, check_extractable=True):
            page_interpreter.process_page(page)
        output_file = open(output_filename, "w")
        output_file.write(string_io.getvalue().replace("\n\n", "\n"))
        input_file.close()
        converter.close()
        output_file.close() 

    def docx_converter(self, input_filename, output_filename):
        input_file = open(input_filename, "rb")
        output_file = open(output_filename, "w", encoding="utf-8")
        output_file.write(docx2txt.process(input_file).replace("\n\n", "\n"))
        input_file.close()
        output_file.close()

    def pptx_converter(self, input_filename, output_filename):
        input_file = open(input_filename, "rb")
        output_file = open(output_filename, "w", encoding="utf-8")
        pptx = Presentation(input_file)
        contents = ""
        for slide in pptx.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        contents += run.text
                contents += "\n"
        output_file.write(contents.replace("\n\n", "\n"))
        input_file.close()
        output_file.close()